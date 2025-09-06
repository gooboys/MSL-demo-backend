from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from pathlib import Path
import os

def iter_shapes_recursive(container):
    """Yield all shapes, descending into groups."""
    for shp in container.shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
            # recurse into group
            for s in iter_shapes_recursive(shp):
                yield s
        else:
            yield shp

def set_text_safe(shape, text, font_name="DM Sans", font_size=10, bold=False):
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        raise ValueError(f"Shape {shape.shape_id} has no text_frame")
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font_name
    f.size = Pt(font_size)
    f.bold = bold

def set_text_by_id(prs, slide_index, shape_id, text, font_name="DM Sans", font_size=10, bold=False):
    slide = prs.slides[slide_index]
    target = None
    for shp in iter_shapes_recursive(slide):
        if shp.shape_id == shape_id:
            target = shp
            break
    if target is None:
        ids = [s.shape_id for s in iter_shapes_recursive(slide)]
        raise RuntimeError(f"Shape id={shape_id} not found on slide index {slide_index}. Present IDs: {ids}")
    set_text_safe(target, text, font_name, font_size, bold)
    # verify write
    got = target.text_frame.text.strip()
    if got != text:
        raise AssertionError(f"Write-check failed. Expected '{text}', saw '{got}'")
    return target

def save_copy(prs, original_path, suffix="_edited"):
    p = Path(original_path)
    out_path = p.with_name(p.stem + suffix + p.suffix)
    prs.save(out_path.as_posix())
    return out_path.as_posix()

# --- use it ---
demo_path = os.path.abspath(os.path.join(os.path.dirname(__file__), "MSL Insight Template.pptx"))
prs = Presentation(demo_path)

# Your 486 is on Slide 2 -> index 1
set_text_by_id(prs, slide_index=1, shape_id=486,
               text="New content here",
               font_name="DM Sans", font_size=20, bold=True)

out = save_copy(prs, demo_path)  # e.g., ".../MSL Insight Template_edited.pptx"
print("Saved:", out)
