from pptx import Presentation
from pptx.slide import Slide, Slides, SlideLayout
from textToElement import set_text
from pathlib import Path
import os

def pptx_maker(input):
  data = input["data"]
  firstSlide = data[0]["output"]
  # print("first slide data:\n", firstSlide)
  secondSlide = data[1]["output"]
  thirdSlide = data[2]["output"]
  return [firstSlide, secondSlide, thirdSlide]

def save_copy(prs, original_path, suffix="_edited"):
  p = Path(original_path)
  out_path = p.with_name(p.stem + suffix + p.suffix)
  prs.save(out_path.as_posix())
  return out_path.as_posix()


demo_path = os.path.join(os.path.dirname(__file__), "MSL Insight Template.pptx")
prs = Presentation(demo_path)
slide = prs.slides[1]
# example: edit shape by id
for shp in slide.shapes:
  if shp.shape_id == 486:
    set_text(shp, "New content here")
    break

copy_path = save_copy(prs, demo_path)  # -> deck_edited.pptx
print("Saved:", copy_path)