from pptx import Presentation
from pptx.oxml.ns import qn
from copy import deepcopy
import os
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE

def table_text(tbl):
  parts = []
  for r in tbl.rows:
    for c in r.cells:
      parts.append(c.text.strip())
  return " | ".join([p for p in parts if p])

def walk_shapes(shapes, indent=0):
  pad = " " * indent
  for sh in shapes:
    sid = getattr(sh, "shape_id", None)
    name = getattr(sh, "name", "")
    stype = sh.shape_type  # MSO_SHAPE_TYPE enum value
    stype_name = MSO_SHAPE_TYPE(stype).name if stype in MSO_SHAPE_TYPE.__members__.values() else str(stype)

    is_ph = getattr(sh, "is_placeholder", False)
    idx = sh.placeholder_format.idx if is_ph else None
    phtype = (str(sh.placeholder_format.type) if is_ph else "-")

    # Extract text when available
    text = ""
    if getattr(sh, "has_text_frame", False):
      text = sh.text.strip()
    elif stype == MSO_SHAPE_TYPE.TABLE:
      text = table_text(sh.table)

    print(f"{pad}id={sid} type={stype_name:<10} name='{name}' "
          f"placeholder_idx={idx} ph_type={phtype} text={text!r}")

    # Recurse into groups
    if stype == MSO_SHAPE_TYPE.GROUP:
      walk_shapes(sh.shapes, indent + 2)

demo_path = os.path.join(os.path.dirname(__file__), "MSL Insight Template.pptx")
demo = Presentation(demo_path)
for si, slide in enumerate(demo.slides):
    print(f"\nSlide {si}")
    walk_shapes(slide.shapes)