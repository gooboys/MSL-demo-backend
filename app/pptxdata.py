from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.shapes.group import GroupShape
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from io import BytesIO
import ast
import os
from typing import Dict, Any, Tuple, List, Set

# ======================
# EMU conversions
# ======================
EMU_PER_INCH = 914400
EMU_PER_CM = 360000
EMU_PER_PT = 12700
EMU_PER_PX = 9525  # assumes 96 DPI

def _to_emu_units(val, units="in"):
  if isinstance(val, (Inches, Pt, Emu)):
    return int(val)
  if isinstance(val, (int, float)):
    u = units.lower()
    if u in ("in", "inch", "inches"):
      return int(val * EMU_PER_INCH)
    if u in ("cm",):
      return int(val * EMU_PER_CM)
    if u in ("pt", "point", "points"):
      return int(val * EMU_PER_PT)
    if u in ("px", "pixel", "pixels"):
      return int(val * EMU_PER_PX)
  raise TypeError("Position/size must be a number with units in {'in','cm','pt','px'} or a pptx unit (Inches/Pt/Emu).")

def _fit_size(nw_emu, nh_emu, max_w_emu, max_h_emu):
  if nw_emu <= 0 or nh_emu <= 0:
    return max_w_emu, max_h_emu
  r = min(max_w_emu / float(nw_emu), max_h_emu / float(nh_emu))
  return int(nw_emu * r), int(nh_emu * r)

# ======================
# Bullet point spacing utils
# ======================

def _add_bullet_spacing(shp, after_pt: int = 4):
  """
  Adds a bit of space *after* each paragraph except the last.
  This visually separates bullets. No effect if there's only one paragraph.
  """
  tf = getattr(shp, "text_frame", None)
  if tf is None:
    return
  paras = list(tf.paragraphs)
  if len(paras) < 2:
    return
  last = len(paras) - 1
  for i, p in enumerate(paras):
    if not (getattr(p, "text", "") or "").strip():
      continue
    # Add spacing between bullets, not after the last one
    p.space_after = Pt(after_pt if i < last else 0)

# ======================
# Color & text utils
# ======================
def hex_to_rgb(hex_code: str) -> tuple[int, int, int]:
  s = hex_code.lstrip('#')
  if len(s) != 6:
    raise ValueError(f"hex must be 6 chars, got '{hex_code}'")
  return tuple(int(s[i:i+2], 16) for i in (0, 2, 4))

def _textify_categories(val) -> str:
    if val is None:
        return ""
    if isinstance(val, (list, tuple, set)):
        return ", ".join(map(str, val))
    if isinstance(val, str):
        # Handle stringified lists: "['Competitive Insights','Education']"
        try:
            parsed = ast.literal_eval(val)
            if isinstance(parsed, (list, tuple, set)):
                return ", ".join(map(str, parsed))
        except (SyntaxError, ValueError):
            pass
        return val
    return str(val)

# ======================
# Shape discovery
# ======================
def _iter_shapes_recursive(container):
  for shp in container.shapes:
    yield shp
    if isinstance(shp, GroupShape):
      for inner in _iter_shapes_recursive(shp):
        yield inner

def _index_shapes_by_name(prs: Presentation):
  """
  Returns:
    by_name: Dict[str, Tuple[int, object]]  -> first occurrence of name
    dups:    Dict[str, List[Tuple[int, object]]] -> all occurrences (when >1)
  """
  by_name: Dict[str, Tuple[int, Any]] = {}
  dups: Dict[str, List[Tuple[int, Any]]] = {}
  for si, slide in enumerate(prs.slides):
    for shp in _iter_shapes_recursive(slide):
      name = getattr(shp, "name", None)
      if not name:
        continue
      if name in by_name:
        dups.setdefault(name, [by_name[name]]).append((si, shp))
      else:
        by_name[name] = (si, shp)
  return by_name, dups

def dump_shape_map(prs: Presentation):
  """
  Debug helper: prints slide -> shape names
  """
  for si, slide in enumerate(prs.slides):
    print(f"[pptx] slide {si}:")
    for shp in _iter_shapes_recursive(slide):
      nm = getattr(shp, "name", None)
      tp = getattr(getattr(shp, "shape_type", None), "name", "?")
      txt = getattr(getattr(shp, "text_frame", None), "text", "")
      if txt:
        txt = txt.replace("\r", " ").replace("\n", " ")
        if len(txt) > 48:
          txt = txt[:45] + "..."
      print(f"  - type={tp:<12} name='{nm}' text='{txt}'")

# ======================
# Text setters (robust)
# ======================
def _set_text_simple(shp, text: str, font="Century Gothic", size=14, color=(40, 36, 111)):
  tf = getattr(shp, "text_frame", None)
  if tf is None:
    raise ValueError(f"no text_frame on {getattr(shp,'name','?')}")
  tf.text = text or ""

  # Style *all* paragraphs and runs
  for p in tf.paragraphs:
    # ensure at least one run exists
    if not p.runs:
      r = p.add_run()
      r.text = ""
    for r in p.runs:
      if font:  r.font.name = font
      if size:  r.font.size = Pt(size)
      if color: r.font.color.rgb = RGBColor(*color)


def _overwrite_shape_text(
  shp,
  text: str,
  font_name: str = "Calibri",
  font_size: int | float = 20,
  font_color: tuple[int, int, int] = (0, 0, 0),
  bold: bool | None = None,
  italic: bool | None = None,
  bullet_gap_pt: int = 4,
):
  tf = getattr(shp, "text_frame", None)
  if tf is None:
    raise ValueError(f"Shape name={getattr(shp, 'name', '?')} has no text frame")

  # 1) Write text as paragraphs (preserves \n)
  tf.text = text or ""
  paras = list(tf.paragraphs)

  # 2) Style all runs uniformly
  for p in paras:
    if not p.runs:
      r = p.add_run(); r.text = ""
    for r in p.runs:
      if font_name:  r.font.name = font_name
      if font_size:  r.font.size = Pt(font_size)
      if font_color: r.font.color.rgb = RGBColor(*font_color)
      if bold is not None:   r.font.bold = bold
      if italic is not None: r.font.italic = italic

  # 3) If multiple paragraphs, try to FORCE bullets via XML and add spacing
  if len(paras) > 1:
    forced_any = False
    for i, p in enumerate(paras):
      # indent level (0 = top level)
      try:
        p.level = 0
      except Exception:
        pass

      # Clear any existing bullet settings, then add buChar •
      try:
        pPr = p._element.get_or_add_pPr()
        # remove buNone/buAutoNum/buChar/buBlip if present
        for tag in ("a:buNone", "a:buAutoNum", "a:buChar", "a:buBlip"):
          el = pPr.find(qn(tag))
          if el is not None:
            pPr.remove(el)
        buChar = OxmlElement("a:buChar")
        buChar.set("char", u"\u2022")
        pPr.append(buChar)

        # Hanging indent so text aligns after the bullet
        # marL = 18pt, indent = -12pt (in EMUs: 1pt = 12700 EMU)
        marL = int(18 * 12700)
        indent = int(-12 * 12700)
        pPr.set("marL", str(marL))
        pPr.set("indent", str(indent))

        forced_any = True
      except Exception:
        # If XML manipulation fails, we'll fall back below.
        pass

      # Space between bullets (except last)
      try:
        p.space_after = Pt(bullet_gap_pt if i < len(paras) - 1 else 0)
      except Exception:
        pass

    # 4) Fallback: if bullets still not showing, prefix a visible glyph and keep spacing
    if not forced_any:
      for i, p in enumerate(paras):
        # skip empty lines
        t = p.text or ""
        if t.strip():
          # add a visible bullet glyph
          if not t.lstrip().startswith("•"):
            p.clear()  # clear runs safely, keep paragraph
            r = p.add_run()
            r.text = f"• {t}"
            if font_name:  r.font.name = font_name
            if font_size:  r.font.size = Pt(font_size)
            if font_color: r.font.color.rgb = RGBColor(*font_color)
            if bold is not None:   r.font.bold = bold
            if italic is not None: r.font.italic = italic
        try:
          p.space_after = Pt(bullet_gap_pt if i < len(paras) - 1 else 0)
        except Exception:
          pass

# ======================
# Image placement
# ======================
def insert_image_fit_units(
  prs,
  slide_idx: int,
  image_bytes: bytes,
  box_w, box_h,
  pos_x, pos_y,
  units: str = "in"
):
  """
  Place an image (bytes) on slide `slide_idx`, scaled to FIT inside a box of (box_w x box_h)
  whose top-left corner is at (pos_x, pos_y), all in the chosen `units`.
  Returns the picture shape.
  """
  max_w_emu = _to_emu_units(box_w, units)
  max_h_emu = _to_emu_units(box_h, units)
  left_emu  = _to_emu_units(pos_x, units)
  top_emu   = _to_emu_units(pos_y, units)

  slide = prs.slides[slide_idx]
  stream = BytesIO(image_bytes)
  pic = slide.shapes.add_picture(stream, left_emu, top_emu)

  native_w = pic.width
  native_h = pic.height
  fit_w, fit_h = _fit_size(native_w, native_h, max_w_emu, max_h_emu)

  pic.width = fit_w
  pic.height = fit_h
  pic.left = left_emu
  pic.top = top_emu
  return pic

# ======================
# Core editor (name-first, slide-agnostic)
# ======================
def editPPTX(pres, ref: Dict[int, Dict], items: Dict[int, Tuple[str, int]], debug: bool = False):
  """
  Edits shapes by NAME; slide indices in `items` are hints only.
  Survives slide reorders and most template edits as long as names remain stable.

  Args:
    pres  : Presentation
    ref   : dict[int, dict]  -> per-id config (contains text/font/etc.)
    items : dict[int, (shape_name, slide_idx_hint)]
    debug : print duplicate & missing name diagnostics
  """
  # headings using simple setter
  special_ids: Set[int] = {}
  special_names = { items[i][0] for i in special_ids if i in items }

  # name-keyed config from (ref, items)
  ref_by_name: Dict[str, Dict] = {}
  for shape_id, cfg in ref.items():
    info = items.get(shape_id)
    if not info:
      continue
    shape_name, _ = info
    # normalize categories fields to text when present
    if "categories" in cfg and "text" not in cfg:
      cfg["text"] = _textify_categories(cfg["categories"])
    ref_by_name[shape_name] = cfg

  by_name, dups = _index_shapes_by_name(pres)

  if debug and dups:
    print("[pptx] Duplicate shape names detected (will use the first occurrence):")
    for nm, locs in dups.items():
      print(f"  - {nm}: " + ", ".join([f"slide {si}" for si, _ in locs]))

  missing: List[str] = []
  for name, cfg in ref_by_name.items():
    entry = by_name.get(name)
    if entry is None:
      missing.append(name)
      continue
    slide_idx, shp = entry
    try:
      if name in special_names:
        _set_text_simple(
          shp,
          text=cfg.get("text", ""),
          font=cfg.get("font", "Calibri"),
          size=cfg.get("font_size", 12),
          color=cfg.get("font_color", (0, 0, 0)),
        )
      else:
        _overwrite_shape_text(
          shp,
          text=cfg.get("text", ""),
          font_name=cfg.get("font", "Calibri"),
          font_size=cfg.get("font_size", 12),
          font_color=cfg.get("font_color", (0, 0, 0)),
          bold=cfg.get("bold", None),
          italic=cfg.get("italic", None),
        )
    except Exception as e:
      if debug:
        print(f"[pptx] Failed writing '{name}' on slide {slide_idx}: {e}")

  if debug and missing:
    print("[pptx] Names not found in current template (removed/renamed?):")
    for nm in missing:
      print(f"  - {nm}")

# ======================
# Public entry
# ======================
def true_replacement(
  stats: Dict[str, Any],
  patient: List[Dict[str, Any]],
  education: List[Dict[str, Any]],
  competitive: List[Dict[str, Any]],
  single: List[Dict[str, Any]],
  template_path: str | None = None,
  debug: bool = False
) -> bytes:
  """
  Loads the template, injects text+graphs, returns PPTX bytes.
  """
  base_dir = os.path.dirname(os.path.abspath(__file__))
  template_path = template_path or os.path.join(base_dir, "New Acquis Template.pptx")

  if not os.path.exists(template_path):
    raise FileNotFoundError(f"Template not found at: {template_path}")

  prs = Presentation(template_path)

  def safe_get(lst, idx, key, default=""):
    return (lst[idx].get(key) if 0 <= idx < len(lst) and isinstance(lst[idx], dict) else default)

  # Category order (kept for reference)
  order = [
    "Access Insights",
    "Patient Management / Care Insights",
    "Clinical Development Insights",
    "Competitive Insights",
    "Product Insights (Drug Science)",
    "Education",
    "Logistics",
    "Adverse Event (AE) Insights",
    "Other"
  ]
  catcount = [stats.get('category_count', {}).get(key, 0) for key in order]

  # Build ref (id -> cfg). NOTE: categories fields normalized to text via editor.
  ref: Dict[int, Dict[str, Any]] = {
    1362: {"text": safe_get(single, 0, "Raw CRM Input (from MSL)"),
          "font": "Century Gothic", "font_size": 10, "font_color": hex_to_rgb("28246f"), "bold": True},
    1364: {"text": safe_get(single, 0, "idea"),
          "font": "Century Gothic", "font_size": 7, "font_color": hex_to_rgb("28246f"), "bold": True},
    1365: {"text": safe_get(single, 1, "idea"),
          "font": "Century Gothic", "font_size": 7, "font_color": hex_to_rgb("28246f"), "bold": True},
    1366: {"text": safe_get(single, 3, "idea"),
          "font": "Century Gothic", "font_size": 7, "font_color": hex_to_rgb("28246f"), "bold": True},
    1367: {"text": safe_get(single, 4, "idea"),
          "font": "Century Gothic", "font_size": 7, "font_color": hex_to_rgb("28246f"), "bold": True},
    1368: {"text": safe_get(single, 2, "idea"),
          "font": "Century Gothic", "font_size": 6, "font_color": hex_to_rgb("28246f"), "bold": True},

    1375: {"text": safe_get(single, 0, "value_classification_rationale"),
          "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},
    1376: {"text": safe_get(single, 1, "value_classification_rationale"),
          "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},
    1377: {"text": safe_get(single, 3, "value_classification_rationale"),
          "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},
    1378: {"text": safe_get(single, 4, "value_classification_rationale"),
          "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},
    1379: {"text": safe_get(single, 2, "value_classification_rationale"),
          "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},

    # categories: accept list or str
    1381: {"text": _textify_categories(safe_get(single, 0, "categories")),
          "font": "Century Gothic", "font_size": 10, "font_color": (255, 255, 255)},
    1382: {"text": _textify_categories(safe_get(single, 1, "categories")),
          "font": "Century Gothic", "font_size": 10, "font_color": (255, 255, 255)},

    1385: {"text": safe_get(single, 0, "categorization_rationale"),
          "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},
    1386: {"text": safe_get(single, 1, "categorization_rationale"),
          "font": "Century Gothic", "font_size": 7, "font_color": hex_to_rgb("28246f")},
    1388: {"text": safe_get(single, 2, "categorization_rationale"),
          "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},

    # Slide 2
    1405: {"text": stats.get("Reporting_Dates", ""),
          "font": "Century Gothic", "font_size": 16, "font_color": hex_to_rgb("28246f")},
    1409: {"text": str(stats.get("deployedMSLS", "")),
          "font": "Century Gothic", "font_size": 16, "font_color": hex_to_rgb("28246f")},
    1411: {"text": "Total: " + str(100),
          "font": "Century Gothic", "font_size": 16, "font_color": hex_to_rgb("28246f")},
    1407: {"text": "\n".join(stats.get('Congresses', [])),
          "font": "Century Gothic", "font_size": 16, "font_color": hex_to_rgb("28246f")},
    1413: {"text": str(132),
          "font": "Century Gothic", "font_size": 16, "font_color": hex_to_rgb("28246f")},
    1434: {"text": str(189),
          "font": "Century Gothic", "font_size": 16, "font_color": hex_to_rgb("28246f")},

    1415: {"text": str(catcount[0]), "font": "Century Gothic", "font_size": 11, "font_color": hex_to_rgb("333333"), "bold":True},
    1417: {"text": str(catcount[1]), "font": "Century Gothic", "font_size": 11, "font_color": hex_to_rgb("333333"), "bold":True},
    1419: {"text": str(catcount[2]), "font": "Century Gothic", "font_size": 11, "font_color": hex_to_rgb("333333"), "bold":True},
    1421: {"text": str(catcount[3]), "font": "Century Gothic", "font_size": 11, "font_color": hex_to_rgb("333333"), "bold":True},
    1423: {"text": str(catcount[4]), "font": "Century Gothic", "font_size": 11, "font_color": hex_to_rgb("333333"), "bold":True},
    1425: {"text": str(catcount[5]), "font": "Century Gothic", "font_size": 11, "font_color": hex_to_rgb("333333"), "bold":True},
    1427: {"text": str(catcount[6]), "font": "Century Gothic", "font_size": 11, "font_color": hex_to_rgb("333333"), "bold":True},
    1429: {"text": str(catcount[7]), "font": "Century Gothic", "font_size": 11, "font_color": hex_to_rgb("333333"), "bold":True},
    1431: {"text": str(catcount[8]), "font": "Century Gothic", "font_size": 11, "font_color": hex_to_rgb("333333"), "bold":True},

    # Slide 3: Patient
    1441: {"text": f"Theme 1 (n={len(patient[0].get('other_sources', [])) + 3})",
          "font": "Century Gothic", "font_size": 14, "font_color": hex_to_rgb("28246f")},
    1450: {"text": f"Theme 2 (n={len(patient[1].get('other_sources', [])) + 3})",
          "font": "Century Gothic", "font_size": 14, "font_color": hex_to_rgb("28246f")},
    1459: {"text": f"Theme 3 (n={len(patient[2].get('other_sources', [])) + 3})",
          "font": "Century Gothic", "font_size": 14, "font_color": hex_to_rgb("28246f")},

    1443: {"text": patient[0].get("gap_definition", ""),
          "font": "Century Gothic", "font_size": 10, "font_color": hex_to_rgb("28246f")},
    1452: {"text": patient[1].get("gap_definition", ""),
          "font": "Century Gothic", "font_size": 10, "font_color": hex_to_rgb("28246f")},
    1461: {"text": patient[2].get("gap_definition", ""),
          "font": "Century Gothic", "font_size": 10, "font_color": hex_to_rgb("28246f")},

    1445: {"text": "\n".join([f"id {q.get('id')}: '{q.get('quote','')}'" for q in patient[0].get("representative_quotes", [])]),
          "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},
    1454: {"text": "\n".join([f"id {q.get('id')}: '{q.get('quote','')}'" for q in patient[1].get("representative_quotes", [])]),
          "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},
    1463: {"text": "\n".join([f"id {q.get('id')}: '{q.get('quote','')}'" for q in patient[2].get("representative_quotes", [])]),
          "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},

    1447: {"text": f"1: {patient[0].get('root_cause_questions', ['',''])[0] or ''}\n"
          f"2: {patient[0].get('root_cause_questions', ['',''])[1] or ''}",
      "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},
    1456: {"text": f"1: {patient[1].get('root_cause_questions', ['',''])[0] or ''}\n"
          f"2: {patient[1].get('root_cause_questions', ['',''])[1] or ''}",
      "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},
    1465: {"text": f"1: {patient[2].get('root_cause_questions', ['',''])[0] or ''}\n"
          f"2: {patient[2].get('root_cause_questions', ['',''])[1] or ''}",
      "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},

    # Slide 4: Education
    1470: {"text": f"Theme 1 (n={len(education[0].get('other_sources', [])) + 3})",
          "font": "Century Gothic", "font_size": 14, "font_color": hex_to_rgb("28246f")},
    1479: {"text": f"Theme 2 (n={len(education[1].get('other_sources', [])) + 3})",
          "font": "Century Gothic", "font_size": 14, "font_color": hex_to_rgb("28246f")},
    1488: {"text": f"Theme 3 (n={len(education[2].get('other_sources', [])) + 3})",
          "font": "Century Gothic", "font_size": 14, "font_color": hex_to_rgb("28246f")},

    1472: {"text": education[0].get("gap_definition", ""),
          "font": "Century Gothic", "font_size": 10, "font_color": hex_to_rgb("28246f")},
    1481: {"text": education[1].get("gap_definition", ""),
          "font": "Century Gothic", "font_size": 10, "font_color": hex_to_rgb("28246f")},
    1490: {"text": education[2].get("gap_definition", ""),
          "font": "Century Gothic", "font_size": 10, "font_color": hex_to_rgb("28246f")},

    1474: {"text": "\n".join([f"id {q.get('id')}: '{q.get('quote','')}'" for q in education[0].get("representative_quotes", [])]),
          "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},
    1483: {"text": "\n".join([f"id {q.get('id')}: '{q.get('quote','')}'" for q in education[1].get("representative_quotes", [])]),
          "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},
    1492: {"text": "\n".join([f"id {q.get('id')}: '{q.get('quote','')}'" for q in education[2].get("representative_quotes", [])]),
          "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},

    1476: {"text": f"1: {education[0].get('root_cause_questions', ['',''])[0] or ''}\n"
          f"2: {education[0].get('root_cause_questions', ['',''])[1] or ''}",
      "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},
    1485: {"text": f"1: {education[1].get('root_cause_questions', ['',''])[0] or ''}\n"
          f"2: {education[1].get('root_cause_questions', ['',''])[1] or ''}",
      "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},
    1494: {"text": f"1: {education[2].get('root_cause_questions', ['',''])[0] or ''}\n"
          f"2: {education[2].get('root_cause_questions', ['',''])[1] or ''}",
      "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},

    # Slide 5: Competitive
    1499: {"text": f"Theme 1 (n={len(competitive[0].get('other_sources', [])) + 3})",
          "font": "Century Gothic", "font_size": 14, "font_color": hex_to_rgb("28246f")},
    1508: {"text": f"Theme 2 (n={len(competitive[1].get('other_sources', [])) + 3})",
          "font": "Century Gothic", "font_size": 14, "font_color": hex_to_rgb("28246f")},
    1517: {"text": f"Theme 3 (n={len(competitive[2].get('other_sources', [])) + 3})",
          "font": "Century Gothic", "font_size": 14, "font_color": hex_to_rgb("28246f")},

    1501: {"text": competitive[0].get("gap_definition", ""),
          "font": "Century Gothic", "font_size": 10, "font_color": hex_to_rgb("28246f")},
    1510: {"text": competitive[1].get("gap_definition", ""),
          "font": "Century Gothic", "font_size": 10, "font_color": hex_to_rgb("28246f")},
    1519: {"text": competitive[2].get("gap_definition", ""),
          "font": "Century Gothic", "font_size": 10, "font_color": hex_to_rgb("28246f")},

    1503: {"text": "\n".join([f"id {q.get('id')}: '{q.get('quote','')}'" for q in competitive[0].get("representative_quotes", [])]),
          "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},
    1512: {"text": "\n".join([f"id {q.get('id')}: '{q.get('quote','')}'" for q in competitive[1].get("representative_quotes", [])]),
          "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},
    1521: {"text": "\n".join([f"id {q.get('id')}: '{q.get('quote','')}'" for q in competitive[2].get("representative_quotes", [])]),
          "font": "Century Gothic", "font_size": 8, "font_color": hex_to_rgb("28246f")},

    1505: {"text": f"1: {competitive[0].get('root_cause_questions', ['',''])[0] or ''}\n"
        f"2: {competitive[0].get('root_cause_questions', ['',''])[1] or ''}",
      "font": "Century Gothic", "font_size": 9, "font_color": hex_to_rgb("28246f")},
    1514: {"text": f"1: {competitive[1].get('root_cause_questions', ['',''])[0] or ''}\n"
        f"2: {competitive[1].get('root_cause_questions', ['',''])[1] or ''}",
      "font": "Century Gothic", "font_size": 9, "font_color": hex_to_rgb("28246f")},
    1523: {"text": f"1: {competitive[2].get('root_cause_questions', ['',''])[0] or ''}\n"
        f"2: {competitive[2].get('root_cause_questions', ['',''])[1] or ''}",
      "font": "Century Gothic", "font_size": 9, "font_color": hex_to_rgb("28246f")}
  }

  print("buggy text: ","\n".join([f"id {q.get('id')}: '{q.get('quote','')}'" for q in competitive[0].get("representative_quotes", [])]))
  print("\n".join([f"id {q.get('id')}: '{q.get('quote','')}'" for q in competitive[1].get("representative_quotes", [])]))
  print("\n".join([f"id {q.get('id')}: '{q.get('quote','')}'" for q in competitive[2].get("representative_quotes", [])]))

  # items: id -> (shape_name, slide_idx_hint).
  items: Dict[int, Tuple[str, int]] = {1350: ('Rectangle: Rounded Corners 136', 4), 1351: ('Google Shape;499;p44', 5), 1352: ('Google Shape;500;p44', 5), 1353: ('Google Shape;502;p44', 5), 1354: ('Google Shape;503;p44', 5), 1355: ('Google Shape;504;p44', 5), 1356: ('Google Shape;506;p44', 5), 1357: ('Google Shape;507;p44', 5), 1358: ('Google Shape;508;p44', 5), 1359: ('Google Shape;510;p44', 5), 1360: ('Google Shape;511;p44', 5), 1361: ('Google Shape;512;p44', 5), 1362: ('Google Shape;513;p44', 5), 1363: ('Google Shape;514;p44', 5), 1364: ('Google Shape;515;p44', 5), 1365: ('Google Shape;516;p44', 5), 1366: ('Google Shape;517;p44', 5), 1367: ('Google Shape;518;p44', 5), 1368: ('Google Shape;519;p44', 5), 1369: ('Google Shape;525;p44', 5), 1370: ('Google Shape;526;p44', 5), 1371: ('Google Shape;527;p44', 5), 1372: ('Google Shape;528;p44', 5), 1373: ('Google Shape;529;p44', 5), 1374: ('Google Shape;530;p44', 5), 1375: ('Google Shape;531;p44', 5), 1376: ('Google Shape;532;p44', 5), 1377: ('Google Shape;533;p44', 5), 1378: ('Google Shape;534;p44', 5), 1379: ('Google Shape;535;p44', 5), 1380: ('Google Shape;536;p44', 5), 1381: ('Google Shape;537;p44', 5), 1382: ('Google Shape;538;p44', 5), 1383: ('Google Shape;539;p44', 5), 1384: ('Google Shape;540;p44', 5), 1385: ('Google Shape;541;p44', 5), 1386: ('Google Shape;542;p44', 5), 1387: ('Google Shape;543;p44', 5), 1388: ('Google Shape;544;p44', 5), 1389: ('Google Shape;545;p44', 5), 1390: ('Google Shape;433;p43', 6), 1391: ('Google Shape;434;p43', 6), 1392: ('Rectangle: Rounded Corners 1', 6), 1393: ('Rectangle: Rounded Corners 6', 6), 1394: ('Rectangle: Rounded Corners 2', 6), 1395: ('Rectangle 8', 6), 1396: ('Google Shape;565;p45', 7), 1397: ('Google Shape;566;p45', 7), 1398: ('Google Shape;567;p45', 7), 1399: ('Google Shape;568;p45', 7), 1400: ('Google Shape;569;p45', 7), 1401: ('Google Shape;570;p45', 7), 1402: ('Google Shape;571;p45', 7), 1403: ('Google Shape;572;p45', 7), 1404: ('Google Shape;573;p45', 7), 1405: ('Google Shape;575;p45', 7), 1406: ('Google Shape;576;p45', 7), 1407: ('Google Shape;578;p45', 7), 1408: ('Google Shape;579;p45', 7), 1409: ('Google Shape;581;p45', 7), 1410: ('Google Shape;582;p45', 7), 1411: ('Google Shape;584;p45', 7), 1412: ('Google Shape;585;p45', 7), 1413: ('Google Shape;587;p45', 7), 1414: ('Google Shape;588;p45', 7), 1415: ('Google Shape;590;p45', 7), 1416: ('Google Shape;591;p45', 7), 1417: ('Google Shape;593;p45', 7), 1418: ('Google Shape;594;p45', 7), 1419: ('Google Shape;596;p45', 7), 1420: ('Google Shape;597;p45', 7), 1421: ('Google Shape;599;p45', 7), 1422: ('Google Shape;600;p45', 7), 1423: ('Google Shape;602;p45', 7), 1424: ('Google Shape;603;p45', 7), 1425: ('Google Shape;605;p45', 7), 1426: ('Google Shape;606;p45', 7), 1427: ('Google Shape;608;p45', 7), 1428: ('Google Shape;609;p45', 7), 1429: ('Google Shape;611;p45', 7), 1430: ('Google Shape;612;p45', 7), 1431: ('Google Shape;614;p45', 7), 1432: ('Google Shape;615;p45', 7), 1433: ('Google Shape;616;p45', 7), 1434: ('Google Shape;587;p45;s8;sid59', 7), 1435: ('Google Shape;588;p45;s8;sid60', 7), 1436: ('Google Shape;616;p45;s8;sid61', 7), 1437: ('Google Shape;622;p46', 8), 1438: ('Google Shape;623;p46', 8), 1439: ('Google Shape;624;p46', 8), 1440: ('Google Shape;625;p46', 8), 1441: ('Google Shape;626;p46', 8), 1442: ('Google Shape;627;p46', 8), 1443: ('Google Shape;628;p46', 8), 1444: ('Google Shape;629;p46', 8), 1445: ('Google Shape;630;p46', 8), 1446: ('Google Shape;631;p46', 8), 1447: ('Google Shape;632;p46', 8), 1448: ('Google Shape;633;p46', 8), 1449: ('Google Shape;634;p46', 8), 1450: ('Google Shape;635;p46', 8), 1451: ('Google Shape;636;p46', 8), 1452: ('Google Shape;637;p46', 8), 1453: ('Google Shape;638;p46', 8), 1454: ('Google Shape;639;p46', 8), 1455: ('Google Shape;640;p46', 8), 1456: ('Google Shape;641;p46', 8), 1457: ('Google Shape;642;p46', 8), 1458: ('Google Shape;643;p46', 8), 1459: ('Google Shape;644;p46', 8), 1460: ('Google Shape;645;p46', 8), 1461: ('Google Shape;646;p46', 8), 1462: ('Google Shape;647;p46', 8), 1463: ('Google Shape;648;p46', 8), 1464: ('Google Shape;649;p46', 8), 1465: ('Google Shape;650;p46', 8), 1466: ('Google Shape;656;p47', 9), 1467: ('Google Shape;657;p47', 9), 1468: ('Google Shape;658;p47', 9), 1469: ('Google Shape;659;p47', 9), 1470: ('Google Shape;660;p47', 9), 1471: ('Google Shape;661;p47', 9), 1472: ('Google Shape;662;p47', 9), 1473: ('Google Shape;663;p47', 9), 1474: ('Google Shape;664;p47', 9), 1475: ('Google Shape;665;p47', 9), 1476: ('Google Shape;666;p47', 9), 1477: ('Google Shape;667;p47', 9), 1478: ('Google Shape;668;p47', 9), 1479: ('Google Shape;669;p47', 9), 1480: ('Google Shape;670;p47', 9), 1481: ('Google Shape;671;p47', 9), 1482: ('Google Shape;672;p47', 9), 1483: ('Google Shape;673;p47', 9), 1484: ('Google Shape;674;p47', 9), 1485: ('Google Shape;675;p47', 9), 1486: ('Google Shape;676;p47', 9), 1487: ('Google Shape;677;p47', 9), 1488: ('Google Shape;678;p47', 9), 1489: ('Google Shape;679;p47', 9), 1490: ('Google Shape;680;p47', 9), 1491: ('Google Shape;681;p47', 9), 1492: ('Google Shape;682;p47', 9), 1493: ('Google Shape;683;p47', 9), 1494: ('Google Shape;684;p47', 9), 1495: ('Google Shape;690;p48', 10), 1496: ('Google Shape;691;p48', 10), 1497: ('Google Shape;692;p48', 10), 1498: ('Google Shape;693;p48', 10), 1499: ('Google Shape;694;p48', 10), 1500: ('Google Shape;695;p48', 10), 1501: ('Google Shape;696;p48', 10), 1502: ('Google Shape;697;p48', 10), 1503: ('Google Shape;698;p48', 10), 1504: ('Google Shape;699;p48', 10), 1505: ('Google Shape;700;p48', 10), 1506: ('Google Shape;701;p48', 10), 1507: ('Google Shape;702;p48', 10), 1508: ('Google Shape;703;p48', 10), 1509: ('Google Shape;704;p48', 10), 1510: ('Google Shape;705;p48', 10), 1511: ('Google Shape;706;p48', 10), 1512: ('Google Shape;707;p48', 10), 1513: ('Google Shape;708;p48', 10), 1514: ('Google Shape;709;p48', 10), 1515: ('Google Shape;710;p48', 10), 1516: ('Google Shape;711;p48', 10), 1517: ('Google Shape;712;p48', 10), 1518: ('Google Shape;713;p48', 10), 1519: ('Google Shape;714;p48', 10), 1520: ('Google Shape;715;p48', 10), 1521: ('Google Shape;716;p48', 10), 1522: ('Google Shape;717;p48', 10), 1523: ('Google Shape;718;p48', 10), 1524: ('Google Shape;295;p39;s12;sid295', 11), 1525: ('Google Shape;296;p39;s12;sid296', 11), 1526: ('Google Shape;297;p39;s12;sid297', 11), 1527: ('Google Shape;298;p39;s12;sid298', 11), 1528: ('Google Shape;299;p39;s12;sid299', 11), 1529: ('Google Shape;300;p39;s12;sid300', 11)}

  # Add graphs (confirm indices if your slide order changed)
  # If slide order is volatile, consider naming a placeholder shape and deriving its slide via _index_shapes_by_name.
  insert_image_fit_units(
    prs,
    slide_idx=7,
    image_bytes=stats['graph1'],
    box_w=4.2,
    box_h=2.4,
    pos_x=3.65,
    pos_y=1.9,
    units="in"
  )

  insert_image_fit_units(
    prs,
    slide_idx=7,
    image_bytes=stats['graph2'],
    box_w=4.3,
    box_h=3,
    pos_x=8.45,
    pos_y=1.9,
    units="in"
  )

  # Populate text by name, regardless of slide moves
  editPPTX(prs, ref, items, debug=debug)

  buf = BytesIO()
  prs.save(buf)
  buf.seek(0)
  return buf.getvalue()
