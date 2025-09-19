from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import BytesIO
import os

# EMU conversions
EMU_PER_INCH = 914400
EMU_PER_CM = 360000
EMU_PER_PT = 12700
EMU_PER_PX = 9525  # assumes 96 DPI

def _to_emu_units(val, units="in"):
  if isinstance(val, (Inches, Pt, Emu)):
    return int(val)
  if isinstance(val, (int, float)):
    u = units.lower()
    if u in ("in", "inch", "inches"):
      return int(val * EMU_PER_INCH)
    if u in ("cm",):
      return int(val * EMU_PER_CM)
    if u in ("pt", "point", "points"):
      return int(val * EMU_PER_PT)
    if u in ("px", "pixel", "pixels"):
      return int(val * EMU_PER_PX)
  raise TypeError("Position/size must be a number with units in {'in','cm','pt','px'} or a pptx unit (Inches/Pt/Emu).")

def _fit_size(nw_emu, nh_emu, max_w_emu, max_h_emu):
  if nw_emu <= 0 or nh_emu <= 0:
    return max_w_emu, max_h_emu
  r = min(max_w_emu / float(nw_emu), max_h_emu / float(nh_emu))
  return int(nw_emu * r), int(nh_emu * r)

def insert_image_fit_units(
  prs,
  slide_idx: int,
  image_bytes: bytes,
  box_w, box_h,           # size of the bounding box
  pos_x, pos_y,           # top-left position of the box
  units: str = "in"       # 'in', 'cm', 'pt', or 'px'
):
  """
  Place an image (bytes) on slide `slide_idx`, scaled to FIT inside a box of (box_w x box_h)
  whose top-left corner is at (pos_x, pos_y), all in the chosen `units`.
  Returns the picture shape.
  """
  # Convert all to EMU
  max_w_emu = _to_emu_units(box_w, units)
  max_h_emu = _to_emu_units(box_h, units)
  left_emu  = _to_emu_units(pos_x, units)
  top_emu   = _to_emu_units(pos_y, units)

  slide = prs.slides[slide_idx]
  stream = BytesIO(image_bytes)

  # Add picture at the intended anchor, then size it
  pic = slide.shapes.add_picture(stream, left_emu, top_emu)

  # Native size (EMU)
  native_w = pic.width
  native_h = pic.height

  # Compute fit size
  fit_w, fit_h = _fit_size(native_w, native_h, max_w_emu, max_h_emu)

  # Apply size and keep anchored to the same top-left
  pic.width = fit_w
  pic.height = fit_h
  pic.left = left_emu
  pic.top = top_emu

  return pic


def _iter_shapes_recursive(shapes, path=""):
  for shp in shapes:
    yield (path, shp)
    if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
      # Recurse with a breadcrumb path for clarity
      new_path = f"{path}/{getattr(shp, 'name', 'Group') or 'Group'}[{shp.shape_id}]"
      for sub in _iter_shapes_recursive(shp.shapes, new_path):
        yield sub

def find_shape_by_id_recursive(slide, shape_id: int):
  for _, shp in _iter_shapes_recursive(slide.shapes):
    if shp.shape_id == shape_id:
      return shp
  return None

def replace_text_by_id(slide, shape_id, new_text,
                       font_name="Calibri", font_size=20, font_color=(0, 0, 0),
                       bold=None, italic=None):
  shp = find_shape_by_id_recursive(slide, shape_id)
  if shp is None:
    raise ValueError(f"No shape found with ID {shape_id} (check slide index and that IDs haven't changed)")

  if not hasattr(shp, "text_frame") or shp.text_frame is None:
    raise ValueError(f"Shape with ID {shape_id} has no text frame (type={shp.shape_type.name})")

  tf = shp.text_frame
  tf.clear()
  p = tf.paragraphs[0]
  run = p.add_run()
  run.text = new_text

  run.font.name = font_name
  run.font.size = Pt(font_size)
  run.font.color.rgb = RGBColor(*font_color)
  if bold is not None:
    run.font.bold = bold
  if italic is not None:
    run.font.italic = italic


def full_replacement(stats, patient, education, competitive):
  base_dir = os.path.dirname(os.path.abspath(__file__))
  template_path = os.path.join(base_dir, "Acquis Template.pptx")
  print('base_dir: '+ base_dir)
  print('template_path: '+template_path)
  prs = Presentation(template_path)
  

  # FILLERRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRR
  # Patient Management Slide
  pmslide = prs.slides[4]  # change if global search shows a different index
  # Theme size:
  supporting1 = str(len(patient[0]["other_sources"])+3)
  replace_text_by_id(
    pmslide,
    shape_id=73,
    new_text="Theme 1 (n="+supporting1+")",
    font_name="Century Gothic Bold",  # family only
    font_size=14,
    font_color=(48, 25, 52),
    italic=False
  )
  supporting2 = str(len(patient[1]["other_sources"])+3)
  replace_text_by_id(
    pmslide,
    shape_id=74,
    new_text="Theme 2 (n="+supporting2+")",
    font_name="Century Gothic Bold",  # family only
    font_size=14,
    font_color=(48, 25, 52),
    italic=False
  )
  supporting3 = str(len(patient[2]["other_sources"])+3)
  replace_text_by_id(
    pmslide,
    shape_id=75,
    new_text="Theme 3 (n="+supporting3+")",
    font_name="Century Gothic Bold",  # family only
    font_size=14,
    font_color=(48, 25, 52),
    italic=False
  )

  # Main Theme Descriptions
  theme1 = patient[0]["gap_definition"]
  replace_text_by_id(
    pmslide,
    shape_id=60,
    new_text=theme1,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(255, 255, 255),
    italic=False
  )
  theme2 = patient[1]["gap_definition"]
  replace_text_by_id(
    pmslide,
    shape_id=64,
    new_text=theme2,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(255, 255, 255),
    italic=False
  )
  theme3 = patient[2]["gap_definition"]
  replace_text_by_id(
    pmslide,
    shape_id=65,
    new_text=theme3,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(255, 255, 255),
    italic=False
  )

  # Quotation Gaps
  gap1 = patient[0]["representative_quotes"]
  formgap1 = "\n".join([f"id {q['id']}: '{q['quote']}'" for q in gap1])
  replace_text_by_id(
    pmslide,
    shape_id=71,
    new_text=formgap1,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )

  gap2 = patient[1]["representative_quotes"]
  formgap2 = "\n".join([f"id {q['id']}: '{q['quote']}'" for q in gap2])
  replace_text_by_id(
    pmslide,
    shape_id=83,
    new_text=formgap2,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )
  # THE ONE BELOW NEEDS FORMATTING
  gap3 = patient[2]["representative_quotes"]
  formgap3 = "\n".join([f"id {q['id']}: '{q['quote']}'" for q in gap3])
  replace_text_by_id(
    pmslide,
    shape_id=92,
    new_text=formgap3,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )

  # Root Causes
  root1 = "1: "+patient[0]["root_cause_questions"][0]+"\n"+"2: "+patient[0]["root_cause_questions"][1]
  replace_text_by_id(
    pmslide,
    shape_id=79,
    new_text=root1,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )
  root2 = "1: "+patient[1]["root_cause_questions"][0]+"\n"+"2: "+patient[1]["root_cause_questions"][1]
  replace_text_by_id(
    pmslide,
    shape_id=88,
    new_text=root2,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )
  root3 = "1: "+patient[2]["root_cause_questions"][0]+"\n"+"2: "+patient[2]["root_cause_questions"][1]
  replace_text_by_id(
    pmslide,
    shape_id=97,
    new_text=root3,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )

  # FILLERRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRR
  # Education Slide
  pmslide = prs.slides[5]  # change if global search shows a different index
  # Theme size:
  supporting1 = str(len(education[0]["other_sources"])+3)
  replace_text_by_id(
    pmslide,
    shape_id=73,
    new_text="Theme 1 (n="+supporting1+")",
    font_name="Century Gothic Bold",  # family only
    font_size=14,
    font_color=(48, 25, 52),
    italic=False
  )
  supporting2 = str(len(education[1]["other_sources"])+3)
  replace_text_by_id(
    pmslide,
    shape_id=74,
    new_text="Theme 2 (n="+supporting2+")",
    font_name="Century Gothic Bold",  # family only
    font_size=14,
    font_color=(48, 25, 52),
    italic=False
  )
  supporting3 = str(len(education[2]["other_sources"])+3)
  replace_text_by_id(
    pmslide,
    shape_id=75,
    new_text="Theme 3 (n="+supporting3+")",
    font_name="Century Gothic Bold",  # family only
    font_size=14,
    font_color=(48, 25, 52),
    italic=False
  )

  # Main Theme Descriptions
  theme1 = education[0]["gap_definition"]
  replace_text_by_id(
    pmslide,
    shape_id=60,
    new_text=theme1,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(255, 255, 255),
    italic=False
  )
  theme2 = education[1]["gap_definition"]
  replace_text_by_id(
    pmslide,
    shape_id=64,
    new_text=theme2,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(255, 255, 255),
    italic=False
  )
  theme3 = education[2]["gap_definition"]
  replace_text_by_id(
    pmslide,
    shape_id=65,
    new_text=theme3,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(255, 255, 255),
    italic=False
  )

  # Quotation Gaps
  gap1 = education[0]["representative_quotes"]
  formgap1 = "\n".join([f"id {q['id']}: '{q['quote']}'" for q in gap1])
  replace_text_by_id(
    pmslide,
    shape_id=71,
    new_text=formgap1,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )

  gap2 = education[1]["representative_quotes"]
  formgap2 = "\n".join([f"id {q['id']}: '{q['quote']}'" for q in gap2])
  replace_text_by_id(
    pmslide,
    shape_id=83,
    new_text=formgap2,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )
  # THE ONE BELOW NEEDS FORMATTING
  gap3 = education[2]["representative_quotes"]
  formgap3 = "\n".join([f"id {q['id']}: '{q['quote']}'" for q in gap3])
  replace_text_by_id(
    pmslide,
    shape_id=92,
    new_text=formgap3,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )

  # Root Causes
  root1 = "1: "+education[0]["root_cause_questions"][0]+"\n"+"2: "+education[0]["root_cause_questions"][1]
  replace_text_by_id(
    pmslide,
    shape_id=79,
    new_text=root1,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )
  root2 = "1: "+education[1]["root_cause_questions"][0]+"\n"+"2: "+education[1]["root_cause_questions"][1]
  replace_text_by_id(
    pmslide,
    shape_id=88,
    new_text=root2,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )
  root3 = "1: "+education[2]["root_cause_questions"][0]+"\n"+"2: "+education[2]["root_cause_questions"][1]
  replace_text_by_id(
    pmslide,
    shape_id=97,
    new_text=root3,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )

  # FILLERRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRR
  # competition Slide
  pmslide = prs.slides[6]  # change if global search shows a different index
  # Theme size:
  supporting1 = str(len(competitive[0]["other_sources"])+3)
  replace_text_by_id(
    pmslide,
    shape_id=73,
    new_text="Theme 1 (n="+supporting1+")",
    font_name="Century Gothic Bold",  # family only
    font_size=14,
    font_color=(48, 25, 52),
    italic=False
  )
  supporting2 = str(len(competitive[1]["other_sources"])+3)
  replace_text_by_id(
    pmslide,
    shape_id=74,
    new_text="Theme 2 (n="+supporting2+")",
    font_name="Century Gothic Bold",  # family only
    font_size=14,
    font_color=(48, 25, 52),
    italic=False
  )
  supporting3 = str(len(competitive[2]["other_sources"])+3)
  replace_text_by_id(
    pmslide,
    shape_id=75,
    new_text="Theme 3 (n="+supporting3+")",
    font_name="Century Gothic Bold",  # family only
    font_size=14,
    font_color=(48, 25, 52),
    italic=False
  )

  # Main Theme Descriptions
  theme1 = competitive[0]["gap_definition"]
  replace_text_by_id(
    pmslide,
    shape_id=60,
    new_text=theme1,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(255, 255, 255),
    italic=False
  )
  theme2 = competitive[1]["gap_definition"]
  replace_text_by_id(
    pmslide,
    shape_id=64,
    new_text=theme2,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(255, 255, 255),
    italic=False
  )
  theme3 = competitive[2]["gap_definition"]
  replace_text_by_id(
    pmslide,
    shape_id=65,
    new_text=theme3,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(255, 255, 255),
    italic=False
  )

  # Quotation Gaps
  gap1 = competitive[0]["representative_quotes"]
  formgap1 = "\n".join([f"id {q['id']}: '{q['quote']}'" for q in gap1])
  replace_text_by_id(
    pmslide,
    shape_id=71,
    new_text=formgap1,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )

  gap2 = competitive[1]["representative_quotes"]
  formgap2 = "\n".join([f"id {q['id']}: '{q['quote']}'" for q in gap2])
  replace_text_by_id(
    pmslide,
    shape_id=83,
    new_text=formgap2,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )
  # THE ONE BELOW NEEDS FORMATTING
  gap3 = competitive[2]["representative_quotes"]
  formgap3 = "\n".join([f"id {q['id']}: '{q['quote']}'" for q in gap3])
  replace_text_by_id(
    pmslide,
    shape_id=92,
    new_text=formgap3,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )

  # Root Causes
  root1 = "1: "+competitive[0]["root_cause_questions"][0]+"\n"+"2: "+competitive[0]["root_cause_questions"][1]
  replace_text_by_id(
    pmslide,
    shape_id=79,
    new_text=root1,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )
  root2 = "1: "+competitive[1]["root_cause_questions"][0]+"\n"+"2: "+competitive[1]["root_cause_questions"][1]
  replace_text_by_id(
    pmslide,
    shape_id=88,
    new_text=root2,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )
  root3 = "1: "+competitive[2]["root_cause_questions"][0]+"\n"+"2: "+competitive[2]["root_cause_questions"][1]
  replace_text_by_id(
    pmslide,
    shape_id=97,
    new_text=root3,
    font_name="Century Gothic",  # family only
    font_size=9,
    font_color=(48, 25, 52),
    italic=False
  )

  # FILLERRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRRR
  # Stats Slide
  pmslide = prs.slides[3]
  Dates = stats['Reporting_Dates']
  replace_text_by_id(
    pmslide,
    shape_id=203,
    new_text=Dates,
    font_name="Century Gothic",  # family only
    font_size=10,
    font_color=(255, 255, 255),
    italic=False
  )

  deployed = stats['deployedMSLS']
  replace_text_by_id(
    pmslide,
    shape_id=238,
    new_text=str(deployed),
    font_name="Century Gothic",  # family only
    font_size=10,
    font_color=(255, 255, 255),
    italic=False
  )

  hcpStats = 'Total: '+str(stats['totalInteractions']) +'\nAcademic Setting HCPs: '+str(stats['AcademicSettings'])+'\nCommunity Setting HCPs: '+str(stats['CommunitySettings'])
  replace_text_by_id(
    pmslide,
    shape_id=276,
    new_text=hcpStats,
    font_name="Century Gothic",  # family only
    font_size=10,
    font_color=(255, 255, 255),
    italic=False
  )

  insightCount = stats['InsightCount']
  replace_text_by_id(
    pmslide,
    shape_id=27,
    new_text=str(insightCount),
    font_name="Century Gothic",  # family only
    font_size=10,
    font_color=(255, 255, 255),
    italic=False
  )

  Congresses = "\n".join(stats['Congresses'])
  replace_text_by_id(
    pmslide,
    shape_id=235,
    new_text=Congresses,
    font_name="Century Gothic",  # family only
    font_size=10,
    font_color=(255, 255, 255),
    italic=False
  )

  # Image Processing
  insert_image_fit_units(
    prs,
    slide_idx=3,
    image_bytes=stats['graph1'],
    box_w=6,
    box_h=4,           # size of the bounding box
    pos_x=4, 
    pos_y=2,           # top-left position of the box
    units="in"       # 'in', 'cm', 'pt', or 'px'
  )

  insert_image_fit_units(
    prs,
    slide_idx=3,
    image_bytes=stats['graph2'],
    box_w=6,
    box_h=4,           # size of the bounding box
    pos_x=8.2,
    pos_y=2,           # top-left position of the box
    units="in"       # 'in', 'cm', 'pt', or 'px'
  )

  # prs.save("out.pptx")
  buf = BytesIO()
  prs.save(buf)
  return buf.getvalue()